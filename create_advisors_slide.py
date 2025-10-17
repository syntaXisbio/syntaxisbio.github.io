#!/usr/bin/env python3
"""
Create a PowerPoint slide with SyntaxisBio advisors information
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

def create_advisors_slide():
    # Create presentation
    prs = Presentation()
    
    # Add slide with blank layout
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SyntaxisBio Advisory Board"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)  # Professional blue
    
    # Advisor data
    advisors = [
        {
            "name": "David Raymes",
            "title": "Advisor",
            "description": "Accomplished entrepreneur, seasoned executive, and startup investor. Comprehensive experience across B2B and high-tech industries with successful M&A track record.",
            "image": "images/DavidRaymes_circular.png"
        },
        {
            "name": "Steve Guengerich", 
            "title": "Advisor",
            "description": "Technology entrepreneur with 30+ years experience. Founding shareholder of BSG Corporation (acquired for $330M) and Senior Advisor for Innovation at UT Dallas.",
            "image": "images/SteveGuengerich_circular.png"
        }
    ]
    
    # Position advisors in the center (2-person layout)
    start_x = Inches(3.5)  # More centered for 2 people
    card_width = Inches(3.5)
    spacing = Inches(4.5)
    start_y = Inches(1.8)
    image_size = Inches(2.2)  # Slightly larger for 2-person layout
    
    for i, advisor in enumerate(advisors):
        x_pos = start_x + (i * spacing)
        
        # Add circular image if file exists
        if os.path.exists(advisor["image"]):
            # Add image
            img = slide.shapes.add_picture(advisor["image"], 
                                         x_pos + (card_width - image_size) / 2, 
                                         start_y, 
                                         image_size, image_size)
            # Image is already circular
        else:
            # Add circular placeholder if no image
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          x_pos + (card_width - image_size) / 2,
                                          start_y,
                                          image_size, image_size)
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray
            circle.line.color.rgb = RGBColor(0, 102, 204)  # Blue border
            circle.line.width = Pt(3)
        
        # Add name below image
        name_y = start_y + image_size + Inches(0.2)
        name_box = slide.shapes.add_textbox(x_pos, name_y, card_width, Inches(0.4))
        name_frame = name_box.text_frame
        name_frame.text = advisor["name"]
        name_para = name_frame.paragraphs[0]
        name_para.alignment = PP_ALIGN.CENTER
        name_para.font.size = Pt(18)
        name_para.font.bold = True
        name_para.font.color.rgb = RGBColor(51, 65, 85)  # Dark gray
        
        # Add title below name
        title_y = name_y + Inches(0.4)
        title_box = slide.shapes.add_textbox(x_pos, title_y, card_width, Inches(0.4))
        title_frame = title_box.text_frame
        title_frame.text = advisor["title"]
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 102, 204)  # Professional blue
        
        # Add description below title
        desc_y = title_y + Inches(0.5)
        desc_box = slide.shapes.add_textbox(x_pos, desc_y, card_width, Inches(2))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.margin_left = Inches(0.1)
        desc_frame.margin_right = Inches(0.1)
        desc_frame.text = advisor["description"]
        desc_para = desc_frame.paragraphs[0]
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = RGBColor(100, 116, 139)  # Light gray
        desc_para.line_spacing = 1.2
        
        # Add card-like background
        card_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x_pos - Inches(0.1), start_y - Inches(0.1),
                                        card_width + Inches(0.2), Inches(4.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = RGBColor(248, 250, 252)  # Very light gray
        card_bg.line.color.rgb = RGBColor(226, 232, 240)  # Light border
        card_bg.line.width = Pt(1)
        # Move background to back
        slide.shapes._element.remove(card_bg._element)
        slide.shapes._element.insert(0, card_bg._element)
    
    # Add company info at bottom
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.6))
    footer_frame = footer_box.text_frame
    footer_frame.text = "SyntaxisBio • Unclonable Cellular Identification • syntaxisbio.com"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = RGBColor(100, 116, 139)
    
    # Save presentation
    prs.save('SyntaxisBio_Advisors_Slide.pptx')
    print("PowerPoint advisors slide created: SyntaxisBio_Advisors_Slide.pptx")

if __name__ == "__main__":
    create_advisors_slide()